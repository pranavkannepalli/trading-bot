"""File/URL parsers. Each returns a ParsedDoc dataclass."""

from __future__ import annotations

import asyncio
import csv
import email
import io
import json
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import httpx

logger = logging.getLogger(__name__)


# ── Data model ────────────────────────────────────────────────────────────────

@dataclass
class ParsedDoc:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)
    source: str = ""


class UnsupportedFileTypeError(ValueError):
    """Raised when the file extension is not handled."""


# ── Helpers ───────────────────────────────────────────────────────────────────

def _strip_rst_directives(text: str) -> str:
    """Remove RST directives (.. foo:: ) and field lists."""
    text = re.sub(r"\.\. \w[\w-]*::.*", "", text)
    text = re.sub(r"^\s*:\w[\w -]*:.*$", "", text, flags=re.MULTILINE)
    return text


def _html_to_text(html_content: str, url: str = "") -> str:
    """Extract readable text from HTML, preferring main content."""
    try:
        from readability import Document
        doc = Document(html_content)
        main_html = doc.summary()
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(main_html, "html.parser")
        return soup.get_text(separator="\n", strip=True)
    except Exception:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html_content, "html.parser")
        # Remove script/style/nav tags
        for tag in soup(["script", "style", "nav", "header", "footer", "aside"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)


# ── Extension dispatch ────────────────────────────────────────────────────────

def parse_file(path: Path) -> ParsedDoc:
    """Synchronous file parser. Call from executor for async contexts."""
    suffix = path.suffix.lower()

    # ── Plain text / Markdown / RST ──────────────────────────────────────────
    if suffix in {".md", ".txt", ".rst", ".text"}:
        text = path.read_text(encoding="utf-8", errors="replace")
        if suffix == ".rst":
            text = _strip_rst_directives(text)
        return ParsedDoc(text=text, source=str(path), metadata={"type": suffix.lstrip(".")})

    # ── PDF ──────────────────────────────────────────────────────────────────
    if suffix == ".pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(path))
            pages_text = []
            for page_num, page in enumerate(doc):
                pages_text.append(f"[Page {page_num + 1}]\n{page.get_text()}")
            doc.close()
            return ParsedDoc(
                text="\n\n".join(pages_text),
                source=str(path),
                metadata={"type": "pdf", "pages": len(pages_text)},
            )
        except ImportError:
            raise UnsupportedFileTypeError("PyMuPDF not installed — cannot parse PDF")

    # ── HTML ─────────────────────────────────────────────────────────────────
    if suffix in {".html", ".htm"}:
        html = path.read_text(encoding="utf-8", errors="replace")
        return ParsedDoc(
            text=_html_to_text(html),
            source=str(path),
            metadata={"type": "html"},
        )

    # ── DOCX ─────────────────────────────────────────────────────────────────
    if suffix == ".docx":
        try:
            from docx import Document
            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return ParsedDoc(
                text="\n\n".join(paragraphs),
                source=str(path),
                metadata={"type": "docx"},
            )
        except ImportError:
            raise UnsupportedFileTypeError("python-docx not installed")

    # ── PPTX ─────────────────────────────────────────────────────────────────
    if suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_parts.append(t)
                if slide_parts:
                    slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(slide_parts))
            return ParsedDoc(
                text="\n\n".join(slides_text),
                source=str(path),
                metadata={"type": "pptx", "slides": len(slides_text)},
            )
        except ImportError:
            raise UnsupportedFileTypeError("python-pptx not installed")

    # ── Excel / CSV ───────────────────────────────────────────────────────────
    if suffix in {".xlsx", ".xls"}:
        try:
            import pandas as pd
            xl = pd.ExcelFile(str(path))
            parts = []
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                parts.append(
                    f"## Sheet: {sheet_name}\n"
                    f"Shape: {df.shape[0]} rows × {df.shape[1]} columns\n"
                    f"Columns: {', '.join(str(c) for c in df.columns)}\n\n"
                    f"First 10 rows:\n{df.head(10).to_markdown(index=False)}\n\n"
                    f"Summary stats:\n{df.describe(include='all').to_markdown()}"
                )
            return ParsedDoc(
                text="\n\n".join(parts),
                source=str(path),
                metadata={"type": "xlsx"},
            )
        except ImportError:
            raise UnsupportedFileTypeError("pandas/openpyxl not installed")

    if suffix == ".csv":
        try:
            import pandas as pd
            df = pd.read_csv(str(path))
            text = (
                f"CSV: {path.name}\n"
                f"Shape: {df.shape[0]} rows × {df.shape[1]} columns\n"
                f"Columns: {', '.join(str(c) for c in df.columns)}\n\n"
                f"First 10 rows:\n{df.head(10).to_markdown(index=False)}\n\n"
                f"Summary stats:\n{df.describe(include='all').to_markdown()}"
            )
            return ParsedDoc(text=text, source=str(path), metadata={"type": "csv"})
        except ImportError:
            # Fallback to stdlib csv
            with open(path, newline="", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                rows = list(reader)
            header = rows[0] if rows else []
            text = f"CSV: {path.name}\nColumns: {', '.join(header)}\nRows: {len(rows) - 1}\n\n"
            text += "\n".join(",".join(r) for r in rows[:11])
            return ParsedDoc(text=text, source=str(path), metadata={"type": "csv"})

    # ── JSON / JSONL ──────────────────────────────────────────────────────────
    if suffix == ".json":
        raw = path.read_text(encoding="utf-8", errors="replace")
        try:
            obj = json.loads(raw)
            if isinstance(obj, list):
                sample = obj[:3]
                summary = f"JSON array with {len(obj)} items.\nSample:\n{json.dumps(sample, indent=2)}"
            elif isinstance(obj, dict):
                keys = list(obj.keys())
                summary = f"JSON object with keys: {', '.join(str(k) for k in keys[:20])}\n\nContent:\n{json.dumps(obj, indent=2)[:4000]}"
            else:
                summary = raw[:4000]
        except json.JSONDecodeError:
            summary = raw[:4000]
        return ParsedDoc(text=summary, source=str(path), metadata={"type": "json"})

    if suffix == ".jsonl":
        lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
        sample = []
        for line in lines[:5]:
            try:
                sample.append(json.loads(line))
            except Exception:
                pass
        text = f"JSONL: {len(lines)} records\nSample (first 5):\n{json.dumps(sample, indent=2)}"
        return ParsedDoc(text=text, source=str(path), metadata={"type": "jsonl"})

    # ── EPUB ─────────────────────────────────────────────────────────────────
    if suffix == ".epub":
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(str(path))
            chapters = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text = soup.get_text(separator="\n", strip=True)
                    if text.strip():
                        chapters.append(text)
            return ParsedDoc(
                text="\n\n---\n\n".join(chapters),
                source=str(path),
                metadata={"type": "epub", "chapters": len(chapters)},
            )
        except ImportError:
            raise UnsupportedFileTypeError("ebooklib not installed")

    # ── Code files ────────────────────────────────────────────────────────────
    code_extensions = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".go": "go",
        ".rs": "rust",
        ".sh": "bash",
        ".bash": "bash",
        ".zsh": "zsh",
        ".rb": "ruby",
        ".java": "java",
        ".c": "c",
        ".cpp": "cpp",
        ".h": "c",
        ".hpp": "cpp",
        ".kt": "kotlin",
        ".swift": "swift",
        ".php": "php",
        ".sql": "sql",
        ".yaml": "yaml",
        ".yml": "yaml",
        ".toml": "toml",
        ".ini": "ini",
        ".cfg": "ini",
    }
    if suffix in code_extensions:
        lang = code_extensions[suffix]
        code = path.read_text(encoding="utf-8", errors="replace")
        return ParsedDoc(
            text=f"```{lang}\n{code}\n```",
            source=str(path),
            metadata={"type": "code", "language": lang},
        )

    # ── SRT / VTT subtitle files ──────────────────────────────────────────────
    if suffix in {".srt", ".vtt"}:
        raw = path.read_text(encoding="utf-8", errors="replace")
        # Strip timestamp lines and sequence numbers
        lines = []
        for line in raw.splitlines():
            line = line.strip()
            # Skip sequence numbers, timestamps, and WEBVTT header
            if re.match(r"^\d+$", line):
                continue
            if re.match(r"[\d:,\. ]+-->", line):
                continue
            if line in {"WEBVTT", ""}:
                continue
            lines.append(line)
        return ParsedDoc(
            text=" ".join(lines),
            source=str(path),
            metadata={"type": suffix.lstrip(".")},
        )

    # ── EML ───────────────────────────────────────────────────────────────────
    if suffix == ".eml":
        raw = path.read_bytes()
        msg = email.message_from_bytes(raw)
        parts = [
            f"From: {msg.get('From', '')}",
            f"To: {msg.get('To', '')}",
            f"Subject: {msg.get('Subject', '')}",
            f"Date: {msg.get('Date', '')}",
            "",
        ]
        # Extract body
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        parts.append(payload.decode("utf-8", errors="replace"))
                        break
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                parts.append(payload.decode("utf-8", errors="replace"))
        return ParsedDoc(
            text="\n".join(parts),
            source=str(path),
            metadata={"type": "eml", "subject": msg.get("Subject", "")},
        )

    raise UnsupportedFileTypeError(
        f"Unsupported file type: {suffix!r}. Cannot parse {path.name}"
    )


async def parse_url(url: str) -> ParsedDoc:
    """Async URL parser — fetches HTML and extracts main content."""
    async with httpx.AsyncClient(
        follow_redirects=True,
        timeout=30.0,
        headers={"User-Agent": "Archivum/1.0 (knowledge-base-bot)"},
    ) as client:
        response = await client.get(url)
        response.raise_for_status()

    content_type = response.headers.get("content-type", "")

    if "text/html" in content_type or "application/xhtml" in content_type:
        text = _html_to_text(response.text, url=url)
        # Try to extract title
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(response.text, "html.parser")
        title_tag = soup.find("title")
        title = title_tag.text.strip() if title_tag else url

        return ParsedDoc(
            text=text,
            source=url,
            metadata={"type": "url", "url": url, "title": title, "content_type": content_type},
        )

    if "application/json" in content_type:
        return ParsedDoc(
            text=response.text[:8000],
            source=url,
            metadata={"type": "json_url", "url": url},
        )

    if "text/plain" in content_type or "text/markdown" in content_type:
        return ParsedDoc(
            text=response.text,
            source=url,
            metadata={"type": "text_url", "url": url},
        )

    # Fallback: try to decode as text
    return ParsedDoc(
        text=response.text[:8000],
        source=url,
        metadata={"type": "unknown_url", "url": url},
    )


async def parse_source(source: str | Path) -> ParsedDoc:
    """Top-level dispatcher: handles Path objects and URL strings."""
    if isinstance(source, str) and (source.startswith("http://") or source.startswith("https://")):
        return await parse_url(source)

    path = Path(source)
    if not path.exists():
        raise FileNotFoundError(f"Source file not found: {path}")

    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, parse_file, path)
